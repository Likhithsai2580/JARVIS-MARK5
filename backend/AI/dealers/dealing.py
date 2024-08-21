from gradio_client import Client
from backend.modules.llms import pure_llama3
from backend.modules.filter import filter_python
# Initialize clients
client_img_gen = Client("stabilityai/stable-diffusion-3-medium")
client_img_analysis = Client("KingNish/OpenGPT-4o")


def img_generate(prompt, negative_prompt, width=1024, height=1024):
    try:
        result = client_img_gen.predict(
            prompt=prompt,
            negative_prompt=negative_prompt,
            seed=0,
            randomize_seed=True,
            width=width,
            height=height,
            guidance_scale=5,
            num_inference_steps=28,
            api_name="/infer"
        )
        return result
    except Exception as e:
        return {"error": str(e)}

def analyze_image_or_video(image_or_video_path, prompt):
    try:
        result = client_img_analysis.predict(
            user_prompt={"text": prompt, "files": [image_or_video_path]},
            api_name="/chat"
        )
        return result
    except Exception as e:
        return {"error": str(e)}

pdf_dealer_prompt="""
You are an advanced language model capable of interacting with PDF files and automating tasks. Your user may want to perform various operations on PDF files. Here's how to handle those requests using Python:

File Selection and Path Retrieval:

When the user selects a PDF file, use pyautogui to simulate pressing Ctrl+Shift+C to copy the file path to the clipboard.
Retrieve the path from the clipboard.
Processing the PDF:

Use the retrieved file path to open and manipulate the PDF using Python libraries like PyMuPDF, PyPDF2, or pdfplumber.
Here's an example of the Python code that demonstrates this process:
'''
python
import pyautogui
import pyperclip
import time
import fitz  # PyMuPDF

# Function to simulate keypresses and retrieve file path
def get_file_path():
    pyautogui.hotkey('ctrl', 'shift', 'c')  # Press Ctrl+Shift+C to copy file path
    time.sleep(1)  # Wait for the clipboard to update
    file_path = pyperclip.paste()  # Retrieve the file path from the clipboard
    return file_path

# Function to process the PDF
def process_pdf(file_path):
    # Open the PDF file
    pdf_document = fitz.open(file_path)
    
    # Example: Print the number of pages in the PDF
    print(f"Number of pages: {pdf_document.page_count}")
    
    # Example: Extract text from the first page
    first_page = pdf_document.load_page(0)
    text = first_page.get_text()
    print(f"Text from the first page:\n{text}")

    # Close the document
    pdf_document.close()

# Main workflow
file_path = get_file_path()
if file_path:
    process_pdf(file_path)
else:
    print("No file path found.")
'''
"""

excel_dealer_prompt = """
You are an advanced language model capable of interacting with Excel files and automating tasks. Your user may want to perform various operations on Excel files. Here's how to handle those requests using Python:

File Selection and Path Retrieval:

When the user selects an Excel file, use pyautogui to simulate pressing Ctrl+Shift+C to copy the file path to the clipboard.
Retrieve the path from the clipboard.
Processing the Excel File:

Use the retrieved file path to open and manipulate the Excel file using Python libraries like pandas and openpyxl.
Here's an example of the Python code that demonstrates this process:
'''python
import pyautogui
import pyperclip
import time
import pandas as pd

# Function to simulate keypresses and retrieve file path
def get_file_path():
    pyautogui.hotkey('ctrl', 'shift', 'c')  # Press Ctrl+Shift+C to copy file path
    time.sleep(1)  # Wait for the clipboard to update
    file_path = pyperclip.paste()  # Retrieve the file path from the clipboard
    return file_path

# Function to process the Excel file
def process_excel(file_path):
    # Read the Excel file
    excel_data = pd.ExcelFile(file_path)
    
    # Example: Print sheet names
    print(f"Sheet names: {excel_data.sheet_names}")
    
    # Example: Load the first sheet and print its content
    df = pd.read_excel(file_path, sheet_name=0)
    print(f"Data from the first sheet:\n{df.head()}")

# Main workflow
file_path = get_file_path()
if file_path:
    process_excel(file_path)
else:
    print("No file path found.")
'''
"""

powerpoint_dealer_prompt="""

'''python
You are an advanced language model capable of interacting with PowerPoint files and automating tasks. Your user may want to perform various operations on PowerPoint files. Here’s how to handle those requests using Python:

File Selection and Path Retrieval:

When the user selects a PowerPoint file, use pyautogui to simulate pressing Ctrl+Shift+C to copy the file path to the clipboard.
Retrieve the path from the clipboard.
Processing the PowerPoint:

Use the retrieved file path to open and manipulate the PowerPoint file using Python libraries like python-pptx.
Here’s an example of the Python code demonstrating this process:
'''python
import pyautogui
import pyperclip
import time
from pptx import Presentation

# Function to simulate keypresses and retrieve file path
def get_file_path():
    pyautogui.hotkey('ctrl', 'shift', 'c')  # Press Ctrl+Shift+C to copy file path
    time.sleep(1)  # Wait for the clipboard to update
    file_path = pyperclip.paste()  # Retrieve the file path from the clipboard
    return file_path

# Function to process the PowerPoint
def process_presentation(file_path):
    # Open the PowerPoint file
    presentation = Presentation(file_path)
    
    # Example: Print the number of slides in the presentation
    print(f"Number of slides: {len(presentation.slides)}")
    
    # Example: Extract text from the first slide
    first_slide = presentation.slides[0]
    for shape in first_slide.shapes:
        if hasattr(shape, "text"):
            print(f"Text from first slide: {shape.text}")

# Main workflow
file_path = get_file_path()
if file_path:
    process_presentation(file_path)
else:
    print("No file path found.")

"""

def pdf_dealer(user_input):
    exec(filter_python(pure_llama3(pdf_dealer_prompt + user_input)))

def excel_dealer(user_input):
    exec(filter_python(pure_llama3(excel_dealer_prompt + user_input)))

def powerpoint_dealer(user_input):
    exec(filter_python(pure_llama3(powerpoint_dealer_prompt + user_input)))